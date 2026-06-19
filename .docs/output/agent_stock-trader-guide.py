"""
Stock Trader 사용자 가이드 — 14 슬라이드 빌드 스크립트
Output: /Users/wise/Developer/tmp/Stock-Trader/.docs/Stock-Trader.pptx
"""
import os
import sys

PLUGIN_ROOT = "/Users/wise/.claude/plugins/cache/axlabs/axlabs-mckinsey-pptx/0.2.0"
sys.path.insert(0, PLUGIN_ROOT)

from dataclasses import replace
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from mckinsey_pptx import PresentationBuilder, DEFAULT_THEME
from mckinsey_pptx.base import (
    blank_slide, add_chrome, add_rect, add_textbox, write_paragraph,
)

# ── Korean theme ──────────────────────────────────────────────────────────────
KO_THEME = replace(
    DEFAULT_THEME,
    typography=replace(DEFAULT_THEME.typography, family="Apple SD Gothic Neo"),
    copyright_text="ⓒ 2026 AX Labs",
)

SCREENSHOTS = "/Users/wise/Developer/tmp/Stock-Trader/.docs/Stock-Trader"
OUTPUT      = "/Users/wise/Developer/tmp/Stock-Trader/.docs/Stock-Trader.pptx"

SECTION_MARKER = "Stock Trader v2.31"

b      = PresentationBuilder(theme=KO_THEME, default_section_marker=SECTION_MARKER)
layout = KO_THEME.layout
typo   = KO_THEME.typography
pal    = KO_THEME.palette

SW = layout.slide_width_in    # 13.33"
SH = layout.slide_height_in   # 7.5"
ML = layout.margin_left_in
MR = layout.margin_right_in
FOOTER_TOP = layout.footer_top_in   # ~6.85"


# ── Helper: screenshot slide ──────────────────────────────────────────────────
def add_screenshot_slide(page_num, title_text, img_filename, bullets):
    """
    Left column (~4.4"): McKinsey chrome + annotation bullets
    Right column (~8.3"): embedded screenshot, top-aligned
    """
    slide = blank_slide(b.prs)
    add_chrome(slide, title=title_text, theme=KO_THEME,
               page_number=page_num, section_marker=SECTION_MARKER)

    BODY_TOP = 1.85
    BODY_H   = FOOTER_TOP - BODY_TOP - 0.15

    LEFT_X = ML + 0.10
    LEFT_W = 4.30

    IMG_X = ML + LEFT_W + 0.25
    IMG_W = SW - IMG_X - MR - 0.10
    IMG_Y = BODY_TOP
    IMG_H = BODY_H

    # --- Bullet column ---
    tb = add_textbox(slide, LEFT_X, BODY_TOP, LEFT_W, BODY_H,
                     anchor=MSO_ANCHOR.TOP)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        is_header = item.startswith("##")
        if is_header:
            text = item[2:].strip()
            p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
            run = p.add_run()
            run.text = text
            run.font.size = Pt(9.5)
            run.font.bold = True
            run.font.color.rgb = pal.bright_blue
            run.font.name = typo.family
            if not first:
                p.space_before = Pt(6)
        else:
            p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(9.0)
            run.font.bold = False
            run.font.color.rgb = pal.text_dark
            run.font.name = typo.family
            p.space_before = Pt(2)
        first = False

    # --- Screenshot ---
    img_path = os.path.join(SCREENSHOTS, img_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(IMG_X), Inches(IMG_Y),
            Inches(IMG_W), Inches(IMG_H),
        )
    else:
        add_rect(slide, IMG_X, IMG_Y, IMG_W, IMG_H,
                 fill=RGBColor(0xEE, 0xEE, 0xEE))
        tb2 = add_textbox(slide, IMG_X + 0.3, IMG_Y + IMG_H / 2 - 0.2,
                          IMG_W - 0.6, 0.4, anchor=MSO_ANCHOR.MIDDLE)
        write_paragraph(tb2.text_frame,
                        f"[스크린샷 없음: {img_filename}]",
                        size=9, color=pal.placeholder_gray,
                        family=typo.family, align=PP_ALIGN.CENTER, first=True)

    b._page = page_num  # keep counter in sync
    return slide


# ═════════════════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ═════════════════════════════════════════════════════════════════════════════
b.add("cover_slide",
      title="Stock Trader 사용자 가이드",
      subtitle="AI 기반 한국 주식시장 트레이딩 대시보드 사용 설명서",
      client="AX Labs",
      date="v2.31  |  2026-06-19",
      confidentiality=None)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 2 — 시스템 개요 (overview_areas: 7 modules)
# ═════════════════════════════════════════════════════════════════════════════
b.add("overview_areas",
      title="시스템 개요 — 7개 핵심 모듈 & 4종 페르소나",
      subtitle="",
      areas=[
          {"name": "대시보드",
           "bullets": ["종목 검색 및 실시간 지수", "현재가·RSI·볼린저 카드", "에이전트 결정 표시"]},
          {"name": "포트폴리오",
           "bullets": ["보유 종목 현황 및 손익", "평가손익·잔여예수금·총기자산", "종목별 수익률·비중"]},
          {"name": "전략·스크리너",
           "bullets": ["9가지 필터 조건 설정", "스캔→매칭 종목 리스트", "결과 클릭→대시보드 이동"]},
          {"name": "리스크 모니터",
           "bullets": ["종목코드 입력 조회", "KRX 시장경보 상태 표시", "공매도 추이 확인"]},
          {"name": "백테스팅",
           "bullets": ["규칙기반 4종 전략", "강화학습 4종 알고리즘", "수익률·샤프·낙폭·승률 결과"]},
          {"name": "에이전트 분석",
           "bullets": ["6단계 자동 파이프라인", "BUY/SELL/HOLD + 신뢰도%", "페르소나별 포지션 비중"]},
          {"name": "서비스 상태",
           "bullets": ["ingest·analysis·agents·risk", "녹색=정상 / 적색=오류", "사이드바 하단 상태 도트"]},
      ],
      call_out="페르소나: 스캘퍼 · 데이 · 스윙 · 포지션")

# ═════════════════════════════════════════════════════════════════════════════
# Slide 3 — 로그인
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=3,
    title_text="로그인",
    img_filename="01-login.png",
    bullets=[
        "## 로그인 방법",
        "이메일 주소 입력",
        "비밀번호 입력 후 로그인 버튼 클릭",
        "## 편의 옵션",
        "아이디 기억하기: 이메일 자동 채움",
        "자동로그인: 재방문 시 자동 인증",
        "## 계정 생성",
        "회원가입 링크 클릭으로 신규 계정 생성",
        "이름·이메일·비밀번호 입력 필요",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 4 — 대시보드 상단
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=4,
    title_text="대시보드 — 상단: TopBar & 핵심 지표 카드",
    img_filename="02-dashboard-top.png",
    bullets=[
        "## TopBar",
        "종목코드/기업명 검색 후 조회 버튼",
        "KOSPI · KOSDAQ · KRX300 실시간 지수",
        "SIMULATION 배지: 모의거래 모드",
        "다크/라이트 모드 토글 (우측)",
        "## 핵심 지표 카드 (5종)",
        "현재가: 실시간 체결가",
        "RSI: 과매수·과매도 판단 지수",
        "볼린저밴드: 상단·하단 밴드값",
        "지표시그널: BUY/SELL/HOLD",
        "에이전트결정: 신뢰도·포지션 비중",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 5 — 대시보드 하단
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=5,
    title_text="대시보드 — 하단: 차트·기술지표·수급·리스크",
    img_filename="02-dashboard-bottom.png",
    bullets=[
        "## 캔들차트 F6.2",
        "분봉 캔들차트 (5분봉 등)",
        "시간대별 평균 수익률 히트맵",
        "요일별 평균 수익률 히트맵",
        "## 기술지표 패널",
        "RSI · MACD · Signal · Hist · ATR",
        "EMA · SMA · VWAP · Close%",
        "## 수급 동향 (Phase A)",
        "기관·외국인·개인 당일 순매수 금액",
        "## 리스크 상태",
        "Stop-Loss / 일일한도 / 포지션한도",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 6 — 포트폴리오
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=6,
    title_text="포트폴리오",
    img_filename="03-portfolio.png",
    bullets=[
        "## 요약 수치 (상단 카드)",
        "보유 종목 수",
        "평가손익 (원)",
        "잔여예수금 (원)",
        "총기자산 (원)",
        "## 보유 포지션 테이블",
        "종목코드 · 종목명 · 보유수량",
        "매수단가 · 현재가 · 평가손익",
        "수익률(%) · 비중(%)",
        "## 종목 선택",
        "행 클릭 → 대시보드 자동 이동",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 7 — 전략·스크리너 필터
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=7,
    title_text="전략·스크리너 — 필터 설정",
    img_filename="04-strategy-screener.png",
    bullets=[
        "## 9가지 필터 조건",
        "마켓: KRX / KOSPI / KOSDAQ",
        "시그널: 전체 / BUY / SELL / HOLD",
        "RSI 범위: 최솟값 ~ 최댓값",
        "최소 거래량",
        "종가 범위: 하한 ~ 상한 (원)",
        "최소 ESG 점수",
        "최대 공매도 비율(%)",
        "결과 수 제한",
        "## 실행",
        "스크리닝 버튼 → 결과 확인",
        "↺백테스팅 링크로 바로 이동 가능",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 8 — 전략·스크리너 결과
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=8,
    title_text="전략·스크리너 — 스크리닝 결과",
    img_filename="04-strategy-screener-results.png",
    bullets=[
        "## 결과 요약",
        "스캔 종목 수 / 매칭 종목 수 표시",
        "## 결과 테이블 컬럼",
        "종목코드 · 종목명",
        "현재가 · 거래량",
        "RSI · 시그널",
        "ESG · 공매도%",
        "## 시그널 색상",
        "BUY = 파랑 / SELL = 빨강 / HOLD = 회색",
        "## 종목 선택",
        "행 클릭 → 해당 종목 대시보드로 이동",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 9 — 리스크 모니터
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=9,
    title_text="리스크 모니터",
    img_filename="05-risk.png",
    bullets=[
        "## 종목 조회",
        "종목코드 입력 후 조회 버튼 클릭",
        "대시보드 선택 종목 자동 연동",
        "## 시장정보 패널",
        "KRX 시장경보 레벨 확인",
        "경보없음(정상) / 경보 1~3단계",
        "## 공매도 추이 패널",
        "일별 공매도 수량 및 비율 통계",
        "데이터 없을 시 '데이터없음' 표시",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 10 — 백테스팅 규칙기반
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=10,
    title_text="백테스팅 — 규칙기반 전략",
    img_filename="06-backtest.png",
    bullets=[
        "## 전략 선택 (4종)",
        "SMA 교차",
        "RSI 임계값",
        "MACD 교차",
        "Q-러닝",
        "## 기간 설정",
        "60 / 180 / 365 / 730일",
        "## 결과 메트릭 (5종)",
        "총수익률 · 샤프지수 · 최대낙폭",
        "승률 · 총거래수",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 11 — 백테스팅 강화학습
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=11,
    title_text="백테스팅 — 강화학습(RL) 전략",
    img_filename="06-backtest-rl-tab.png",
    bullets=[
        "## 알고리즘 (4종)",
        "DQN (Deep Q-Network)",
        "PPO (Proximal Policy Opt.)",
        "A2C (Advantage Actor-Critic)",
        "QR-DQN (Quantile Regression)",
        "## 파라미터 설정",
        "기간: 60/180/365/730일",
        "에피소드 수 (기본 50회)",
        "## 주의사항",
        "학습 시 수 분 소요될 수 있음",
        "완료 후 5종 결과 메트릭 자동 표시",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 12 — 에이전트 분석
# ═════════════════════════════════════════════════════════════════════════════
add_screenshot_slide(
    page_num=12,
    title_text="에이전트 분석 — 6단계 자동 파이프라인",
    img_filename="07-agents.png",
    bullets=[
        "## 실행 방법",
        "종목코드·페르소나 확인",
        "에이전트 실행 버튼 클릭",
        "## 6단계 파이프라인",
        "01 시세 수집 완료",
        "02 기술지표·예측 해석",
        "03 페르소나 비중 산정",
        "04 수급 데이터 처리 (없으면 degrade)",
        "05 시장정보 레벨 확인",
        "06 최종 시그널 결정",
        "## 최종 출력",
        "시그널 · 신뢰도% · 포지션비중% · 근거",
    ],
)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 13 — 사이드바 & 서비스 상태 (three_trends_table)
# ═════════════════════════════════════════════════════════════════════════════
b.add("three_trends_table",
      title="사이드바 · 페르소나 · 서비스 상태",
      trends=[
          {
              "name": "네비게이션 & 페르소나",
              "description": [
                  "사이드바: 6개 메뉴 (대시보드/포트폴리오/전략/리스크/백테스팅/에이전트)",
                  "스캘퍼: 초단타 1~5분 스캘핑 전략",
                  "데이: 당일 청산 일봉 전략",
                  "스윙: 수일~수주 중기 전략",
                  "포지션: 장기 추세 추종 전략",
              ],
              "examples": [
                  "페르소나 전환 즉시 에이전트 비중 재산정",
                  "활성 페르소나 강조 표시",
              ],
          },
          {
              "name": "서비스 상태 도트",
              "description": [
                  "ingest: 시세 수집 서비스",
                  "analysis: 기술지표 분석 엔진",
                  "agents: 에이전트 의사결정",
                  "risk: 리스크 모니터링",
              ],
              "examples": [
                  "녹색 = 정상 작동",
                  "적색 = 연결불가·장애",
                  "회색 = 상태 확인중",
              ],
          },
          {
              "name": "기타 기능",
              "description": [
                  "TopBar ☀/🌙 버튼: 다크/라이트 모드 토글",
                  "사이드바 하단: 로그인 사용자 확인",
                  "로그아웃: 사이드바 하단 버튼 클릭",
              ],
              "examples": [
                  "모든 페이지 동일한 사이드바 사용",
                  "종목 변경: TopBar 검색창으로",
              ],
          },
      ])

# ═════════════════════════════════════════════════════════════════════════════
# Slide 14 — 주요 사용 시나리오 (overview_areas: 4 scenarios)
# overview_areas chosen over phases_table_4 — no hardcoded [TIMELINE] artifacts
# ═════════════════════════════════════════════════════════════════════════════
b.add("overview_areas",
      title="주요 사용 시나리오",
      subtitle="",
      areas=[
          {
              "name": "A. 종목 분석",
              "bullets": [
                  "TopBar에서 종목코드 검색 후 조회",
                  "현재가·RSI·볼린저 지표 카드 확인",
                  "에이전트 실행 → BUY/SELL/HOLD",
                  "신뢰도·포지션 비중·근거 검토",
              ],
          },
          {
              "name": "B. 종목 탐색",
              "bullets": [
                  "전략·스크리너 필터 설정",
                  "스크리닝 실행 → 매칭 종목 확인",
                  "결과 행 클릭 → 대시보드 이동",
                  "에이전트 실행으로 즉시 판단",
              ],
          },
          {
              "name": "C. 전략 검증",
              "bullets": [
                  "백테스팅 규칙기반 전략 선택·실행",
                  "수익률·샤프지수·최대낙폭 확인",
                  "RL 탭 전환 → DQN/PPO 비교",
                  "최적 전략·기간 조합 선정",
              ],
          },
          {
              "name": "D. 리스크 점검",
              "bullets": [
                  "포트폴리오 보유 종목 손익 확인",
                  "Stop-Loss·일일한도·포지션한도 점검",
                  "리스크 모니터 종목코드 조회",
                  "시장 경보·공매도 추이 이상 확인",
              ],
          },
      ])

# ═════════════════════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════════════════════
b.save(OUTPUT)
print(f"Saved: {OUTPUT}")
